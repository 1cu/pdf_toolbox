"""PPTX manipulation actions."""

from __future__ import annotations

from pathlib import Path
from typing import Literal, cast

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pdf_toolbox.actions import action
from pdf_toolbox.actions.pdf_images import QualityChoice, resolve_image_settings
from pdf_toolbox.paths import validate_path
from pdf_toolbox.renderers.pptx import require_pptx_renderer
from pdf_toolbox.utils import logger


def _parse_range_spec(spec: str, total: int) -> list[int]:
    """Return slide numbers described by ``spec`` preserving order."""
    if not spec:
        return list(range(1, total + 1))
    result: list[int] = []
    for raw_part in spec.split(","):
        part = raw_part.strip().replace("n", str(total))
        if not part:
            continue
        if "-" in part:
            start_s, end_s = part.split("-", 1)
            start = int(start_s) if start_s else 1
            end = int(end_s) if end_s else total
            if start < 1 or end > total or end < start:
                msg = f"invalid range {part}"
                raise ValueError(msg)
            result.extend(range(start, end + 1))
        else:
            page = int(part)
            if page < 1 or page > total:
                msg = f"page {page} out of range"
                raise ValueError(msg)
            result.append(page)
    return result


@action(category="PPTX")
def extract_pptx_images(input_pptx: str, out_dir: str | None = None) -> str:
    """Extract embedded images from ``input_pptx``.

    Args:
        input_pptx: Path to the presentation.
        out_dir: Destination directory; defaults to ``<input>_images``.

    Returns:
        Path to the directory containing the extracted images.

    Example:
        >>> extract_pptx_images("slides.pptx")
        'slides_images'
    """
    pres = Presentation(str(validate_path(input_pptx, must_exist=True)))
    base_dir = Path(input_pptx).parent
    target = Path(out_dir) if out_dir else base_dir / f"{Path(input_pptx).stem}_images"
    target = validate_path(target)
    target.mkdir(parents=True, exist_ok=True)

    count = 0
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                filename = target / f"image_{count + 1}.{image.ext}"
                filename.write_bytes(image.blob)
                count += 1
    logger.info("Extracted %d image(s) from %s", count, input_pptx)
    return str(target)


@action(category="PPTX")
def reorder_pptx(
    input_pptx: str,
    range_spec: str,
    output_path: str | None = None,
) -> str:
    """Reorder slides of ``input_pptx`` according to ``range_spec``.

    Args:
        input_pptx: Presentation to reorder.
        range_spec: Slide specification such as ``"1-3,7,4-5,8-n"``.
        output_path: Optional path for the reordered file. Defaults to
            ``<input>_reordered.pptx``.
    """
    pres = Presentation(str(validate_path(input_pptx, must_exist=True)))
    numbers = _parse_range_spec(range_spec, len(pres.slides))
    sld_id_lst = pres.slides._sldIdLst
    slides = list(sld_id_lst)
    sld_id_lst.clear()
    for num in numbers:
        sld_id_lst.append(slides[num - 1])
    out = (
        Path(output_path)
        if output_path
        else Path(input_pptx).with_name(f"{Path(input_pptx).stem}_reordered.pptx")
    )
    out = validate_path(out)
    pres.save(str(out))
    logger.info("Reordered %s to %s", input_pptx, out)
    return str(out)


@action(category="PPTX")
def pptx_to_images(  # noqa: PLR0913  # pdf-toolbox: action interface requires many parameters | issue:-
    input_pptx: str,
    out_dir: str | None = None,
    max_size_mb: float | None = None,
    image_format: Literal["PNG", "JPEG", "TIFF"] = "JPEG",
    quality: int | QualityChoice = "High (95)",
    width: int | None = None,
    height: int | None = None,
    pages: str | None = None,
) -> str:
    """Render a PPTX file to images using the configured provider.

    Args:
        input_pptx: Presentation to render.
        out_dir: Destination directory for exported images.
        max_size_mb: Optional size limit per image in megabytes.
        image_format: Output image format (``"JPEG"``, ``"PNG"``, or ``"TIFF"``).
        quality: JPEG/WebP quality (ignored for other formats).
        width: Optional pixel width; requires ``height``.
        height: Optional pixel height; requires ``width``.
        pages: Optional slide selection using the same syntax as PDF exports.

    Returns:
        Path to the directory containing the images.
    """
    fmt, quality_val, _dpi = resolve_image_settings(
        image_format,
        quality,
        allowed_formats={"PNG", "JPEG", "TIFF"},
    )
    fmt_literal = cast(Literal["PNG", "JPEG", "TIFF"], fmt)
    renderer = require_pptx_renderer()
    return renderer.to_images(
        input_pptx,
        out_dir=out_dir,
        max_size_mb=max_size_mb,
        image_format=fmt_literal,
        quality=quality_val,
        width=width,
        height=height,
        range_spec=pages,
    )


@action(category="PPTX")
def pptx_to_pdf(
    input_pptx: str, output_path: str | None = None, pages: str | None = None
) -> str:
    """Render a PPTX file to PDF using the configured provider."""
    renderer = require_pptx_renderer()
    return renderer.to_pdf(input_pptx, output_path=output_path, range_spec=pages)


__all__ = [
    "extract_pptx_images",
    "pptx_to_images",
    "pptx_to_pdf",
    "reorder_pptx",
]
