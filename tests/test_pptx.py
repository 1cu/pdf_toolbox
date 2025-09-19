"""Tests for PPTX actions."""

from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx import Presentation

from pdf_toolbox import config
from pdf_toolbox.actions import pptx as pptx_actions
from pdf_toolbox.actions.pptx import pptx_to_images, pptx_to_pdf
from pdf_toolbox.renderers import pptx
from pdf_toolbox.renderers import registry as pptx_registry
from pdf_toolbox.renderers.pptx import (
    BasePptxRenderer,
    PptxProviderUnavailableError,
    get_pptx_renderer,
)


@pytest.fixture
def simple_pptx(tmp_path) -> str:
    prs = Presentation()
    for idx in range(5):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Slide {idx + 1}"
    path = tmp_path / "simple.pptx"
    prs.save(path)
    return str(path)


def test_rendering_actions_raise(simple_pptx):
    with pytest.raises(pptx_registry.RendererSelectionError):
        pptx_to_images(simple_pptx)
    with pytest.raises(PptxProviderUnavailableError):
        pptx_to_pdf(simple_pptx)


def test_null_renderer_methods_raise():
    renderer = pptx.NullRenderer()

    with pytest.raises(PptxProviderUnavailableError) as excinfo:
        renderer.to_images("deck.pptx")

    assert excinfo.value.docs_url == pptx.PPTX_PROVIDER_DOCS_URL

    with pytest.raises(PptxProviderUnavailableError):
        renderer.to_pdf("deck.pptx")


def test_ensure_registered_skips_unknown(monkeypatch):
    seen: list[str] = []

    def fake_select(name: str) -> BasePptxRenderer | None:
        seen.append(name)
        return None

    monkeypatch.setattr(pptx, "registry_select", fake_select)

    assert pptx._load_via_registry("missing") is None
    assert seen == ["missing"]


def test_load_via_registry_handles_empty_and_missing(monkeypatch):
    seen: list[str] = []

    def fake_select(name: str) -> BasePptxRenderer | None:
        seen.append(name)
        return None

    monkeypatch.setattr(pptx, "registry_select", fake_select)

    assert pptx._load_via_registry("") is None
    assert pptx._load_via_registry("custom") is None
    assert seen == ["custom"]


def test_get_pptx_renderer_falls_back(monkeypatch, tmp_path):
    cfg_path = tmp_path / "pptx.json"
    cfg_path.write_text(json.dumps({"pptx_renderer": "custom"}))
    monkeypatch.setattr(config, "CONFIG_PATH", cfg_path)
    monkeypatch.setattr(pptx, "_load_via_registry", lambda _name: None)

    renderer = pptx.get_pptx_renderer()

    assert isinstance(renderer, pptx.NullRenderer)


def test_renderer_config(monkeypatch, tmp_path, simple_pptx):
    captured_pdf: dict[str, str | None] = {}
    captured_images: dict[str, object] = {}

    class DummyRenderer(BasePptxRenderer):
        name = "dummy"

        def to_images(  # noqa: PLR0913  # pdf-toolbox: renderer API requires many parameters | issue:-
            self,
            _input_pptx: str,
            out_dir: str | None = None,
            max_size_mb: float | None = None,
            image_format: str = "JPEG",
            quality: int | None = None,
            width: int | None = None,
            height: int | None = None,
            range_spec: str | None = None,
        ) -> str:
            del out_dir, max_size_mb, image_format, quality, width, height, range_spec
            return "ok"

        def to_pdf(
            self,
            _input_pptx: str,
            output_path: str | None = None,
            notes: bool = False,
            handout: bool = False,
            range_spec: str | None = None,
        ) -> str:
            del notes, handout
            target = Path(output_path) if output_path else tmp_path / "fallback.pdf"
            target.write_text("pdf")
            captured_pdf["range_spec"] = range_spec
            return str(target)

    monkeypatch.setattr(
        pptx_registry,
        "_REGISTRY",
        dict(pptx_registry._REGISTRY),
    )
    pptx_registry.register(DummyRenderer)
    monkeypatch.setattr(
        pptx,
        "_load_via_registry",
        lambda name: DummyRenderer() if name == "dummy" else None,
    )
    cfg_path = tmp_path / "cfg.json"
    cfg_path.write_text(json.dumps({"pptx_renderer": "dummy"}))
    monkeypatch.setattr(config, "CONFIG_PATH", cfg_path)

    out_dir = tmp_path / "images"

    def fake_pdf_to_images(
        pdf_path: str,
        *,
        pages: str | None = None,
        image_format: str,
        quality: int,
        **kwargs: object,
    ) -> list[str]:
        captured_images["pdf_path"] = pdf_path
        captured_images["pages"] = pages
        captured_images["image_format"] = image_format
        captured_images["quality"] = quality
        out_dir_obj = kwargs.get("out_dir")
        out_dir_str = str(out_dir_obj) if out_dir_obj is not None else None
        captured_images["out_dir"] = out_dir_str
        captured_images["width"] = kwargs.get("width")
        captured_images["height"] = kwargs.get("height")
        out_dir_path = (
            Path(out_dir_str) if out_dir_str is not None else Path(tmp_path, "images")
        )
        out_dir_path.mkdir(parents=True, exist_ok=True)
        image_path = out_dir_path / "slide-001.jpeg"
        image_path.write_text("image")
        return [str(image_path)]

    monkeypatch.setattr(pptx_actions, "pdf_to_images", fake_pdf_to_images)

    renderer = get_pptx_renderer()
    assert isinstance(renderer, DummyRenderer)
    result_dir = pptx_to_images(simple_pptx, pages="1-2", out_dir=str(out_dir))
    assert Path(result_dir) == Path(out_dir)
    assert captured_images["pages"] == "1-2"
    assert captured_pdf["range_spec"] is None
    assert pptx_to_pdf(simple_pptx, pages="2-3") == str(tmp_path / "fallback.pdf")
    assert captured_pdf["range_spec"] == "2-3"


def test_pptx_to_images_normalises_params(monkeypatch, simple_pptx, tmp_path):
    captured: dict[str, object] = {}

    class DummyRenderer(BasePptxRenderer):
        name = "dummy"

        def to_images(  # noqa: PLR0913  # pdf-toolbox: renderer API requires many parameters | issue:-
            self,
            _input_pptx: str,
            out_dir: str | None = None,
            max_size_mb: float | None = None,
            image_format: str = "JPEG",
            quality: int | None = None,
            width: int | None = None,
            height: int | None = None,
            range_spec: str | None = None,
        ) -> str:
            del out_dir, max_size_mb, image_format, quality, width, height, range_spec
            return "ok"

        def to_pdf(
            self,
            _input_pptx: str,
            output_path: str | None = None,
            notes: bool = False,
            handout: bool = False,
            range_spec: str | None = None,
        ) -> str:
            del notes, handout, range_spec
            target = Path(output_path) if output_path else tmp_path / "slides.pdf"
            target.write_text("pdf")
            return str(target)

    monkeypatch.setattr(
        pptx_registry,
        "_REGISTRY",
        dict(pptx_registry._REGISTRY),
    )
    pptx_registry.register(DummyRenderer)
    monkeypatch.setattr(
        pptx,
        "_load_via_registry",
        lambda name: DummyRenderer() if name == "dummy" else None,
    )
    cfg_path = tmp_path / "cfg.json"
    cfg_path.write_text(json.dumps({"pptx_renderer": "dummy"}))
    monkeypatch.setattr(config, "CONFIG_PATH", cfg_path)

    def fake_pdf_to_images(
        pdf_path: str,
        *,
        pages: str | None = None,
        image_format: str,
        quality: int,
        **kwargs: object,
    ) -> list[str]:
        del pdf_path
        captured["format"] = image_format
        captured["quality"] = quality
        captured["range_spec"] = pages
        out_dir_obj = kwargs.get("out_dir")
        out_dir_str = str(out_dir_obj) if out_dir_obj is not None else None
        captured["out_dir"] = out_dir_str
        captured["width"] = kwargs.get("width")
        captured["height"] = kwargs.get("height")
        output_dir = tmp_path / "rendered"
        output_dir.mkdir(parents=True, exist_ok=True)
        image_path = output_dir / f"slide-001.{image_format.lower()}"
        image_path.write_text("img")
        return [str(image_path)]

    monkeypatch.setattr(pptx_actions, "pdf_to_images", fake_pdf_to_images)

    result_dir = pptx_to_images(
        simple_pptx,
        image_format="png",
        quality="Low (70)",
        pages="1-2",
    )
    assert Path(result_dir) == tmp_path / "rendered"
    assert captured["format"] == "PNG"
    assert captured["quality"] == 70
    assert captured["range_spec"] == "1-2"
    assert captured["out_dir"] == str(Path(simple_pptx).resolve().parent)
    assert captured["width"] is None
    assert captured["height"] is None


def test_pptx_to_images_returns_out_dir_when_empty(monkeypatch, simple_pptx, tmp_path):
    captured: dict[str, object] = {}

    class DummyRenderer(BasePptxRenderer):
        name = "dummy"

        def to_images(  # noqa: PLR0913  # pdf-toolbox: renderer API requires many parameters | issue:-
            self,
            _input_pptx: str,
            out_dir: str | None = None,
            max_size_mb: float | None = None,
            image_format: str = "JPEG",
            quality: int | None = None,
            width: int | None = None,
            height: int | None = None,
            range_spec: str | None = None,
        ) -> str:
            del (
                out_dir,
                max_size_mb,
                image_format,
                quality,
                width,
                height,
                range_spec,
            )
            return "unused"

        def to_pdf(
            self,
            _input_pptx: str,
            output_path: str | None = None,
            notes: bool = False,
            handout: bool = False,
            range_spec: str | None = None,
        ) -> str:
            del notes, handout, range_spec
            target = (
                Path(output_path) if output_path else tmp_path / "fallback-empty.pdf"
            )
            target.write_text("pdf")
            return str(target)

    monkeypatch.setattr(
        pptx_registry,
        "_REGISTRY",
        dict(pptx_registry._REGISTRY),
    )
    pptx_registry.register(DummyRenderer)
    monkeypatch.setattr(
        pptx,
        "_load_via_registry",
        lambda name: DummyRenderer() if name == "dummy" else None,
    )
    cfg_path = tmp_path / "cfg.json"
    cfg_path.write_text(json.dumps({"pptx_renderer": "dummy"}))
    monkeypatch.setattr(config, "CONFIG_PATH", cfg_path)

    def fake_pdf_to_images(pdf_path: str, **_kwargs: object) -> list[str]:
        captured["pdf_path"] = pdf_path
        captured["invoked"] = True
        return []

    monkeypatch.setattr(pptx_actions, "pdf_to_images", fake_pdf_to_images)

    out_dir = tmp_path / "empty"
    result = pptx_to_images(simple_pptx, out_dir=str(out_dir))

    assert result == str(out_dir)
    assert captured["invoked"] is True
    assert Path(captured["pdf_path"]).suffix == ".pdf"


def test_pptx_to_images_returns_temp_dir_when_empty(monkeypatch, simple_pptx, tmp_path):
    captured: dict[str, object] = {}

    class DummyRenderer(BasePptxRenderer):
        name = "dummy"

        def to_images(  # noqa: PLR0913  # pdf-toolbox: renderer API requires many parameters | issue:-
            self,
            _input_pptx: str,
            out_dir: str | None = None,
            max_size_mb: float | None = None,
            image_format: str = "JPEG",
            quality: int | None = None,
            width: int | None = None,
            height: int | None = None,
            range_spec: str | None = None,
        ) -> str:
            del (
                out_dir,
                max_size_mb,
                image_format,
                quality,
                width,
                height,
                range_spec,
            )
            return "unused"

        def to_pdf(
            self,
            _input_pptx: str,
            output_path: str | None = None,
            notes: bool = False,
            handout: bool = False,
            range_spec: str | None = None,
        ) -> str:
            del notes, handout, range_spec
            target = (
                Path(output_path) if output_path else tmp_path / "fallback-temp.pdf"
            )
            target.write_text("pdf")
            return str(target)

    monkeypatch.setattr(
        pptx_registry,
        "_REGISTRY",
        dict(pptx_registry._REGISTRY),
    )
    pptx_registry.register(DummyRenderer)
    monkeypatch.setattr(
        pptx,
        "_load_via_registry",
        lambda name: DummyRenderer() if name == "dummy" else None,
    )
    cfg_path = tmp_path / "cfg.json"
    cfg_path.write_text(json.dumps({"pptx_renderer": "dummy"}))
    monkeypatch.setattr(config, "CONFIG_PATH", cfg_path)

    def fake_pdf_to_images(pdf_path: str, **kwargs: object) -> list[str]:
        captured["pdf_path"] = pdf_path
        captured["out_dir"] = kwargs.get("out_dir")
        return []

    monkeypatch.setattr(pptx_actions, "pdf_to_images", fake_pdf_to_images)

    result = pptx_to_images(simple_pptx)

    expected_dir = Path(simple_pptx).resolve().parent
    assert Path(result) == expected_dir
    assert Path(str(captured["out_dir"])) == expected_dir
    assert Path(captured["pdf_path"]).parent != expected_dir


def test_convert_pptx_to_pdf_cleans_up(monkeypatch, simple_pptx, tmp_path):
    class DummyRenderer(BasePptxRenderer):
        name = "dummy"

        def to_images(  # noqa: PLR0913  # pdf-toolbox: renderer API requires many parameters | issue:-
            self,
            _input_pptx: str,
            out_dir: str | None = None,
            max_size_mb: float | None = None,
            image_format: str = "JPEG",
            quality: int | None = None,
            width: int | None = None,
            height: int | None = None,
            range_spec: str | None = None,
        ) -> str:
            del (
                out_dir,
                max_size_mb,
                image_format,
                quality,
                width,
                height,
                range_spec,
            )
            return "unused"

        def to_pdf(
            self,
            input_pptx: str,
            output_path: str | None = None,
            notes: bool = False,
            handout: bool = False,
            range_spec: str | None = None,
        ) -> str:
            del input_pptx, notes, handout, range_spec
            target = Path(output_path) if output_path else tmp_path / "fallback.pdf"
            target.write_text("pdf")
            return str(target)

    monkeypatch.setattr(pptx_registry, "_REGISTRY", {})
    pptx_registry.register(DummyRenderer)
    monkeypatch.setattr(
        pptx_registry, "get_pptx_renderer_choice", lambda _cfg=None: "dummy"
    )

    with pptx_registry.convert_pptx_to_pdf(simple_pptx) as pdf_path:
        path_obj = Path(pdf_path)
        assert path_obj.exists()
        assert path_obj.suffix == ".pdf"
        tmp_dir = path_obj.parent
        assert tmp_dir.name.startswith("pdf-toolbox-pptx-")
    assert not path_obj.exists()
    assert not tmp_dir.exists()
