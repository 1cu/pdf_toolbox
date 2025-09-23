import json
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

import fitz
from pdf_toolbox import utils


@pytest.fixture(scope="session")
def sample_pdf(tmp_path_factory: pytest.TempPathFactory) -> str:
    base_dir = tmp_path_factory.mktemp("sample-pdf")
    pdf_path = base_dir / "sample.pdf"
    if not pdf_path.exists():
        document = fitz.open()
        try:
            for page_index in range(3):
                page = document.new_page(width=200, height=200)
                page.insert_text((72, 72), f"Page {page_index + 1}")
            document.save(pdf_path)
        finally:
            document.close()
    return str(pdf_path)


@pytest.fixture(scope="session")
def pdf_with_image(tmp_path_factory: pytest.TempPathFactory) -> str:
    base_dir = tmp_path_factory.mktemp("pdf-with-image")
    pdf_path = base_dir / "with_image.pdf"

    img_path = base_dir / "img.png"
    Image.new("RGB", (10, 10), color=(255, 0, 0)).save(img_path)
    doc = fitz.open()
    try:
        page = doc.new_page(width=200, height=200)
        rect = fitz.Rect(0, 0, 10, 10)
        page.insert_text((72, 72), "Hi")
        page.insert_image(rect, filename=str(img_path))
        doc.save(pdf_path)
    finally:
        doc.close()
    with pdf_path.open("ab") as fh:
        fh.write(b"% pad" + b"0" * 1000)
    return str(pdf_path)


@pytest.fixture(autouse=True)
def author_config(tmp_path, monkeypatch):
    config = tmp_path / "pdf_toolbox_config.json"
    config.write_text(json.dumps({"author": "Tester", "email": "tester@example.com"}))
    monkeypatch.setattr(utils, "CONFIG_FILE", config)


@pytest.fixture
def simple_pptx(tmp_path: Path) -> str:
    prs = Presentation()
    for idx in range(5):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Slide {idx + 1}"
    path = tmp_path / "simple.pptx"
    prs.save(str(path))
    return str(path)
