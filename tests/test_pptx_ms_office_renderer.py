from __future__ import annotations

import sys
from pathlib import Path

import pytest

pytest.importorskip("pptx")
from pptx import Presentation

from pdf_toolbox.renderers.ms_office import PptxMsOfficeRenderer, PptxRenderingError

if not sys.platform.startswith("win"):
    pytest.skip("MS Office integration requires Windows", allow_module_level=True)


@pytest.fixture(scope="module")
def renderer() -> PptxMsOfficeRenderer:
    instance = PptxMsOfficeRenderer()
    try:
        instance._require_env()
    except PptxRenderingError as exc:  # pragma: no cover - depends on Windows env  # pdf-toolbox: requires Windows PowerPoint | issue:-
        pytest.skip(f"PowerPoint environment unavailable: {exc}")

    try:
        app = instance._open_app()
    except Exception as exc:  # pragma: no cover - depends on Windows env  # pdf-toolbox: requires Windows PowerPoint | issue:-
        pytest.skip(f"Cannot start PowerPoint automation: {exc}")
    else:
        instance._close_app(app)
    return instance


def _make_pptx(tmp: Path, slides: int = 1) -> Path:
    prs = Presentation()
    for i in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box = slide.shapes.add_textbox(100, 100, 300, 50)
        box.text_frame.text = f"Slide {i + 1}"
    file = tmp / "demo.pptx"
    prs.save(str(file))
    return file


def test_to_pdf(renderer: PptxMsOfficeRenderer, tmp_path: Path) -> None:
    src = _make_pptx(tmp_path, slides=3)
    out = tmp_path / "out.pdf"
    try:
        got = renderer.to_pdf(str(src), str(out), range_spec="1-2")
    except PptxRenderingError as exc:  # pragma: no cover - depends on Windows env  # pdf-toolbox: requires Windows PowerPoint | issue:-
        pytest.skip(f"PowerPoint export unavailable: {exc}")
    assert Path(got) == out
    assert out.exists()
    assert out.stat().st_size > 0


def test_to_images(renderer: PptxMsOfficeRenderer, tmp_path: Path) -> None:
    src = _make_pptx(tmp_path, slides=2)
    out = tmp_path / "img"
    try:
        got = renderer.to_images(
            str(src),
            out_dir=str(out),
            image_format="PNG",
            width=1920,
            height=1080,
        )
    except PptxRenderingError as exc:  # pragma: no cover - depends on Windows env  # pdf-toolbox: requires Windows PowerPoint | issue:-
        pytest.skip(f"PowerPoint export unavailable: {exc}")
    files = list(Path(got).glob("*.png"))
    assert files, "Kein Bild exportiert"
    assert len(files) == 2
    assert all(path.stat().st_size > 0 for path in files)
